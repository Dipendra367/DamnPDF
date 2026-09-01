import asyncio
import shutil
import uuid
from pathlib import Path

from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import FileResponse
from pdf2image import convert_from_path
from pptx import Presentation

from core.config import MAX_FILE_SIZE
from core.limits import conversion_semaphore

router = APIRouter()

TEMP_DIR = Path("temp_files")
TEMP_DIR.mkdir(exist_ok=True)


PAGE_BATCH = 20  # pages rendered per poppler pass, bounds peak memory


def build_pptx(input_path: Path, output_path: Path, img_dir: Path, total_pages: int):
    """Blocking worker: renders each PDF page and assembles a real PPTX.

    Slides are full-page images (the same approach commercial converters
    use) — LibreOffice's native PDF->PPTX export silently yields an empty
    presentation, so we don't use it.
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide_size_set = False

    for start in range(1, total_pages + 1, PAGE_BATCH):
        pages = convert_from_path(
            input_path, dpi=150, first_page=start, last_page=min(start + PAGE_BATCH - 1, total_pages)
        )
        for offset, page in enumerate(pages):
            num = start + offset
            img_path = img_dir / f"page_{num}.jpg"
            page.save(img_path, "JPEG", quality=90)

            if not slide_size_set:
                # match the PDF's page proportions (150 dpi, 914400 EMU per inch)
                prs.slide_width = int(page.width * 914400 / 150)
                prs.slide_height = int(page.height * 914400 / 150)
                slide_size_set = True

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height
            )

    if not slide_size_set:
        raise HTTPException(status_code=500, detail="No pages found in PDF")

    prs.save(output_path)


@router.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    if file.content_type != "application/pdf":
        raise HTTPException(status_code=400, detail="File must be a PDF")

    content = await file.read()

    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Max size is 100MB.")

    job_id = uuid.uuid4().hex[:8]
    input_path = TEMP_DIR / f"{job_id}_input.pdf"
    output_path = TEMP_DIR / f"{job_id}_output.pptx"
    img_dir = TEMP_DIR / f"{job_id}_imgs"
    img_dir.mkdir(exist_ok=True)
    input_path.write_bytes(content)

    from pypdf import PdfReader
    total_pages = len(PdfReader(str(input_path)).pages)

    try:
        async with conversion_semaphore:
            await asyncio.to_thread(build_pptx, input_path, output_path, img_dir, total_pages)

        return FileResponse(
            path=output_path,
            filename="converted.pptx",
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    finally:
        input_path.unlink(missing_ok=True)
        shutil.rmtree(img_dir, ignore_errors=True)
